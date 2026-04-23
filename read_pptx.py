from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

path = "/home/cuongdang/Downloads/CSE472 - Final project.pptx"
prs = Presentation(path)

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1}")
    print('='*60)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  {text}")
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  [IMAGE: {shape.name}]")
