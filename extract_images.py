from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

path = "/home/cuongdang/Downloads/CSE472 - Final project.pptx"
out  = "/home/cuongdang/Development/projects/msu/CloudAndFog/their_imgs"
os.makedirs(out, exist_ok=True)

prs = Presentation(path)
for si, slide in enumerate(prs.slides):
    for shi, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img  = shape.image
            ext  = img.ext
            data = img.blob
            fname = f"s{si+1}_{shi}_{shape.name.replace(' ','_').replace(';','_')}.{ext}"
            fpath = os.path.join(out, fname)
            with open(fpath, "wb") as f:
                f.write(data)
            print(f"Saved: {fname}  ({len(data)//1024} KB)  pos=({shape.left//914400:.2f}in, {shape.top//914400:.2f}in) size=({shape.width//914400:.2f}x{shape.height//914400:.2f}in)")
