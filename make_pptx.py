from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

W, H = Inches(13.33), Inches(7.5)

BG       = RGBColor(0x07, 0x11, 0x1f)
ACCENT   = RGBColor(0x6b, 0x7a, 0xff)
GREEN    = RGBColor(0x3e, 0xcf, 0x8e)
ORANGE   = RGBColor(0xf6, 0xa6, 0x23)
RED      = RGBColor(0xff, 0x6b, 0x6b)
WHITE    = RGBColor(0xff, 0xff, 0xff)
BODY     = RGBColor(0xc8, 0xcd, 0xe8)
MUTED    = RGBColor(0x5a, 0x63, 0x82)
TAG_COL  = RGBColor(0x6b, 0x7a, 0xff)
CODE_COL = RGBColor(0xa8, 0xd8, 0xff)
CODE_BG  = RGBColor(0x0a, 0x0f, 0x1e)

IMG = "their_imgs"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def add_slide():
    sl = prs.slides.add_slide(blank)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def txb(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)


def tag_label(slide, text, x=Inches(0.8), y=Inches(0.32)):
    tb = txb(slide, x, y, Inches(8), Inches(0.38))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = TAG_COL


def heading(slide, text, x=Inches(0.8), y=Inches(0.76), w=Inches(11.5), size=34, color=WHITE):
    tb = txb(slide, x, y, w, Inches(0.85))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return tb


def body_box(slide, x, y, w, h):
    tb = txb(slide, x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bullet(tf, text, size=17, color=BODY, ac_color=ACCENT, bold_words=None):
    p = tf.add_paragraph()
    p.space_before = Pt(5)
    r0 = p.add_run()
    r0.text = "▸  "
    r0.font.size = Pt(size)
    r0.font.color.rgb = ac_color

    if bold_words:
        remaining = text
        for bw in bold_words:
            if bw in remaining:
                parts = remaining.split(bw, 1)
                if parts[0]:
                    r = p.add_run(); r.text = parts[0]; r.font.size = Pt(size); r.font.color.rgb = color
                rb = p.add_run(); rb.text = bw; rb.font.size = Pt(size); rb.font.color.rgb = ac_color; rb.font.bold = True
                remaining = parts[1]
        if remaining:
            r = p.add_run(); r.text = remaining; r.font.size = Pt(size); r.font.color.rgb = color
    else:
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.color.rgb = color


def rect(slide, x, y, w, h, rgb):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb
    sh.line.fill.background()
    return sh


def code_block(slide, x, y, w, h, lines, accent=ACCENT):
    rect(slide, x, y, w, h, CODE_BG)
    rect(slide, x, y, Inches(0.05), h, accent)
    tb = txb(slide, x + Inches(0.12), y + Inches(0.1), w - Inches(0.2), h - Inches(0.2))
    tb.text_frame.word_wrap = False
    first = True
    for line in lines:
        p = tb.text_frame.paragraphs[0] if first else tb.text_frame.add_paragraph()
        first = False
        r = p.add_run(); r.text = line
        r.font.size = Pt(13); r.font.color.rgb = CODE_COL; r.font.name = "Courier New"


def add_img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)


def source_line(slide, text, y=Inches(6.9)):
    tb = txb(slide, Inches(0.8), y, Inches(11.7), Inches(0.38))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.color.rgb = MUTED; r.font.italic = True


def pill_row(slide, labels, x=Inches(0.8), y=Inches(6.55), accent=ACCENT):
    px = x
    for lbl in labels:
        w = Inches(len(lbl) * 0.13 + 0.5)
        tb2 = slide.shapes.add_textbox(px, y, w, Inches(0.34))
        tb2.fill.solid(); tb2.fill.fore_color.rgb = RGBColor(0x10, 0x18, 0x40)
        tb2.line.color.rgb = accent; tb2.line.width = Pt(0.8)
        p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = lbl
        r2.font.size = Pt(11); r2.font.color.rgb = accent
        px += w + Inches(0.12)


# ── SLIDE 1: Title ──────────────────────────────────────────
sl = add_slide()
tag_label(sl, "CSE 472 — Final Project")
tb = txb(sl, Inches(0.8), Inches(1.4), Inches(7.5), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Heterogeneous Volume"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE

p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "Rendering"
r2.font.size = Pt(44); r2.font.bold = True; r2.font.color.rgb = WHITE

p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "using Ray Marching"
r3.font.size = Pt(44); r3.font.bold = True; r3.font.color.rgb = ACCENT

pill_row(sl, ["Raymarching", "Beer-Lambert", "Henyey-Greenstein", "Unity URP", "HLSL"],
         Inches(0.8), Inches(4.7))

add_img(sl, f"{IMG}/s1_1_Google_Shape_55_p13.png",
        Inches(7.5), Inches(0.3), Inches(5.5), Inches(6.8))


# ── SLIDE 2: Modeling ───────────────────────────────────────
sl = add_slide()
tag_label(sl, "Modeling")
heading(sl, "3D Texture from Blender")

tf = body_box(sl, Inches(0.8), Inches(1.72), Inches(4.8), Inches(4.6))
bullet(tf, "Render a volumetric cloud in Blender", 16)
bullet(tf, "Take a screenshot of every Z-slice of the rendered cloud", 16,
       bold_words=["every Z-slice"])
bullet(tf, "Pack all 64 slices into a 64×64×64 3D texture sprite sheet", 16,
       bold_words=["64×64×64"])
bullet(tf, "Import into Unity as a 3D texture — used as the density field for raymarching", 16,
       bold_words=["density field"])

add_img(sl, f"{IMG}/s2_4_Google_Shape_64_p14.png",  # Blender scene
        Inches(0.8), Inches(5.0), Inches(4.0), Inches(2.0))
add_img(sl, f"{IMG}/s2_2_Google_Shape_62_p14.png",  # texture sheet
        Inches(5.1), Inches(1.1), Inches(7.9), Inches(5.9))

source_line(sl, "Source: Want to get Blender clouds into #unity? — https://www.youtube.com/watch?v=IGx_I4nzfxQ")


# ── SLIDE 3: Rendering — Raymarching + IGN ─────────────────
sl = add_slide()
tag_label(sl, "Rendering")
heading(sl, "Raymarching & Artifact Reduction")

tf = body_box(sl, Inches(0.8), Inches(1.72), Inches(4.9), Inches(3.2))
bullet(tf, "Raymarch against the 3D density texture generated from Blender", 16,
       bold_words=["Raymarch"])
bullet(tf, "Cast a ray per pixel; accumulate density and light at fixed step intervals", 16)
bullet(tf, "Interleaved Gradient Noise (IGN) randomizes each step size, hiding banding artifacts from uniform stepping", 16,
       bold_words=["Interleaved Gradient Noise (IGN)"])

# three images from their slide
add_img(sl, f"{IMG}/s3_3_Google_Shape_78_p15.png",   # noise pattern
        Inches(0.8), Inches(5.0), Inches(2.0), Inches(1.9))
add_img(sl, f"{IMG}/s3_4_Google_Shape_79_p15.png",   # artifact image
        Inches(3.0), Inches(5.0), Inches(2.0), Inches(1.9))
add_img(sl, f"{IMG}/s3_2_Google_Shape_77_p15.png",   # clean result
        Inches(5.2), Inches(1.72), Inches(7.7), Inches(4.8))

tb_c = txb(sl, Inches(0.8), Inches(4.8), Inches(4.4), Inches(0.28))
p_c = tb_c.text_frame.paragraphs[0]
r_c = p_c.add_run(); r_c.text = "← Artifact (uniform step)          ← IGN pattern"
r_c.font.size = Pt(9); r_c.font.color.rgb = MUTED

source_line(sl, "Source: https://blog.demofox.org/2022/01/01/interleaved-gradient-noise-a-different-kind-of-low-discrepancy-sequence/")


# ── SLIDE 4: Rendering — Henyey-Greenstein ─────────────────
sl = add_slide()
tag_label(sl, "Rendering")
heading(sl, "Henyey-Greenstein Phase Function", color=WHITE)

tf = body_box(sl, Inches(0.8), Inches(1.72), Inches(11.5), Inches(2.0))
bullet(tf, "Simulates directional light scattering through cloud water droplets", 17,
       bold_words=["directional light scattering"])
bullet(tf, "g > 0 → forward scattering (bright toward the sun)   |   g < 0 → back scattering   |   g = 0 → isotropic", 16,
       ac_color=GREEN)
bullet(tf, "Used with each raymarch step to weight how much light reaches the camera", 16)

add_img(sl, f"{IMG}/s4_3_Google_Shape_89_p16.png",   # HG diagram/graph
        Inches(0.8), Inches(3.9), Inches(5.2), Inches(2.9))
add_img(sl, f"{IMG}/s4_2_Google_Shape_88_p16.png",   # HG formula strip
        Inches(6.2), Inches(3.9), Inches(6.7), Inches(1.4))

code_block(sl, Inches(6.2), Inches(5.5), Inches(6.7), Inches(1.75),
    ["float henyey_greenstein(float cosTheta, float g) {",
     "    float denom = 1.0 + g*g - 2.0*g*cosTheta;",
     "    return (1.0 - g*g) / (4.0*PI * pow(denom, 1.5));",
     "}"], accent=GREEN)

source_line(sl, "Source: https://omlc.org/classroom/ece532/class3/hg.html")


# ── SLIDE 5: Rendering — Secondary Raymarch + Beer-Lambert ─
sl = add_slide()
tag_label(sl, "Rendering")
heading(sl, "Secondary Raymarch & Beer-Lambert")

tf = body_box(sl, Inches(0.8), Inches(1.72), Inches(6.4), Inches(4.0))
bullet(tf, "Secondary raymarch from each sample point toward the main light source", 17,
       bold_words=["Secondary raymarch"])
bullet(tf, "Currently supports 1 directional light source", 17,
       bold_words=["1 directional light"])
bullet(tf, "Beer-Lambert Law calculates light falloff through the cloud volume:", 17,
       bold_words=["Beer-Lambert Law"])

code_block(sl, Inches(0.8), Inches(4.0), Inches(6.2), Inches(1.55),
    ["// Extinction along view ray:",
     "transmittance *= exp(-density * _StepSize);",
     "",
     "// Light shadow along light ray:",
     "lightAtten = exp(-shadowDensity * lightStep * _ShadowDensity);"])

add_img(sl, f"{IMG}/s5_2_Google_Shape_97_p17.png",
        Inches(7.1), Inches(1.3), Inches(5.9), Inches(5.5))


# ── SLIDE 6: Rendering — Final Color ───────────────────────
sl = add_slide()
tag_label(sl, "Rendering")
heading(sl, "Final Color Accumulation")

tf = body_box(sl, Inches(5.0), Inches(1.72), Inches(7.9), Inches(4.0))
bullet(tf, "Final color = sum of all light reaching the camera that is not blocked by the cloud itself", 17,
       bold_words=["not blocked by the cloud"])
bullet(tf, "At each step: inscattered light × phase function × shadow attenuation × transmittance", 17)
bullet(tf, "Alpha = 1 − transmittance — fully opaque cloud = zero transmittance", 17,
       ac_color=ORANGE)
bullet(tf, "Result blended over the scene colour via alpha compositing", 17)

code_block(sl, Inches(5.0), Inches(5.0), Inches(7.9), Inches(1.55),
    ["float3 inscatter = lightColor * phase * density * shadowAtten;",
     "accum  += inscatter * transmittance * (1-exp(-density*step));",
     "transmittance *= exp(-density * step);",
     "float alpha = saturate(1.0 - transmittance);"], accent=ORANGE)

add_img(sl, f"{IMG}/s6_2_Google_Shape_104_p18.gif",
        Inches(0.3), Inches(1.1), Inches(4.4), Inches(5.8))


# ── SLIDE 7: Bonus — Fog Rendering ─────────────────────────
sl = add_slide()
tag_label(sl, "Bonus")
heading(sl, "Fog Rendering")

tf = body_box(sl, Inches(0.8), Inches(1.72), Inches(11.5), Inches(1.8))
bullet(tf, "Reused Henyey-Greenstein and raymarching algorithm to render atmospheric fog", 17,
       bold_words=["Henyey-Greenstein"])
bullet(tf, "Two variants: fullscreen post-process blit and localized box fog", 17,
       bold_words=["fullscreen", "box fog"])

# left image — box fog
add_img(sl, f"{IMG}/s7_2_Google_Shape_111_p19.png",
        Inches(0.4), Inches(3.2), Inches(5.8), Inches(3.9))
# right image — fullscreen fog
add_img(sl, f"{IMG}/s7_3_Google_Shape_112_p19.png",
        Inches(6.5), Inches(3.2), Inches(6.5), Inches(3.9))

tb_l = txb(sl, Inches(0.4), Inches(3.0), Inches(5.8), Inches(0.28))
p_l = tb_l.text_frame.paragraphs[0]; p_l.alignment = PP_ALIGN.CENTER
r_l = p_l.add_run(); r_l.text = "Box of fog"
r_l.font.size = Pt(12); r_l.font.color.rgb = MUTED

tb_r = txb(sl, Inches(6.5), Inches(3.0), Inches(6.5), Inches(0.28))
p_r = tb_r.text_frame.paragraphs[0]; p_r.alignment = PP_ALIGN.CENTER
r_r = p_r.add_run(); r_r.text = "Fullscreen raymarching"
r_r.font.size = Pt(12); r_r.font.color.rgb = MUTED


# ── SLIDE 8: Resources ──────────────────────────────────────
sl = add_slide()
tag_label(sl, "Resources & Materials")
heading(sl, "References")

refs = [
    ("[1]", "Blender Volumetric Cloud — BlenderKit",
     "https://www.blenderkit.com/asset-gallery-detail/0bf13dfc-5729-4b9c-9238-94782c3b7286/?query=category_subtree:weather"),
    ("[2]", "Your First Volumetric Fog Shader | Unity URP — YouTube",
     "https://www.youtube.com/watch?v=8P338C9vYEE"),
    ("[3]", "Want to get Blender clouds into #unity? Try these 9 steps #gamedev — YouTube",
     "https://www.youtube.com/watch?v=IGx_I4nzfxQ"),
    ("[4]", "Henyey-Greenstein Phase Function — Oregon Medical Laser Center",
     "https://omlc.org/classroom/ece532/class3/hg.html"),
    ("[5]", "Interleaved Gradient Noise — blog.demofox.org",
     "https://blog.demofox.org/2022/01/01/interleaved-gradient-noise-a-different-kind-of-low-discrepancy-sequence/"),
]

y = Inches(1.72)
for num, title, url in refs:
    tb_n = txb(sl, Inches(0.8), y, Inches(11.5), Inches(0.28))
    p_n = tb_n.text_frame.paragraphs[0]
    r_n = p_n.add_run(); r_n.text = num
    r_n.font.size = Pt(9); r_n.font.bold = True; r_n.font.color.rgb = ACCENT

    tb_t = txb(sl, Inches(1.3), y + Inches(0.25), Inches(11.0), Inches(0.35))
    p_t = tb_t.text_frame.paragraphs[0]
    r_t = p_t.add_run(); r_t.text = title
    r_t.font.size = Pt(15); r_t.font.color.rgb = WHITE

    tb_u = txb(sl, Inches(1.3), y + Inches(0.58), Inches(11.0), Inches(0.3))
    p_u = tb_u.text_frame.paragraphs[0]
    r_u = p_u.add_run(); r_u.text = url
    r_u.font.size = Pt(11); r_u.font.color.rgb = MUTED; r_u.font.name = "Courier New"

    y += Inches(1.12)


prs.save("presentation.pptx")
print(f"Saved presentation.pptx  ({prs.slides.__len__()} slides)")
